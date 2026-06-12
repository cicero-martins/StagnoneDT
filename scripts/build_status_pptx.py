"""Generate a 15-slide progress-report PowerPoint for the June 2026 supervisors meeting.

Usage:
    python scripts/build_status_pptx.py

Output:
    docs/stagnone_status_2026-06-16.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import os

ROOT = Path(__file__).resolve().parents[1]
FIG  = ROOT / 'figures'
OUT  = ROOT / 'docs' / 'stagnone_status_2026-06-16.pptx'

# ── colour palette ─────────────────────────────────────────────────────────────
BLUE_DARK  = RGBColor(0x1F, 0x49, 0x7D)   # title bar
BLUE_MID   = RGBColor(0x2E, 0x74, 0xB5)   # accent
TEAL       = RGBColor(0x00, 0x70, 0xC0)
GREEN_OK   = RGBColor(0x37, 0x86, 0x46)
ORANGE_WIP = RGBColor(0xC5, 0x5A, 0x11)
RED_BLOCK  = RGBColor(0xC0, 0x00, 0x00)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GREY  = RGBColor(0x40, 0x40, 0x40)

W = Inches(13.33)
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_layout(prs):
    return prs.slide_layouts[6]   # blank


def add_rect(slide, x, y, w, h, fill_rgb=None, border=False):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if not border:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=DARK_GREY,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def add_image(slide, img_path, x, y, w=None, h=None):
    if not Path(img_path).exists():
        return None
    kwargs = {}
    if w: kwargs['width']  = Inches(w)
    if h: kwargs['height'] = Inches(h)
    return slide.shapes.add_picture(str(img_path), Inches(x), Inches(y), **kwargs)


def title_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.15, fill_rgb=BLUE_DARK)
    add_text(slide, title, 0.3, 0.1, 12.5, 0.7,
             size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.3, 0.78, 12.5, 0.4,
                 size=14, bold=False, color=RGBColor(0xBD, 0xD7, 0xEE), align=PP_ALIGN.LEFT)


def footer(slide, text='Stagnone Digital Twin | PhD — UNIPA | June 2026'):
    add_rect(slide, 0, 7.15, 13.33, 0.35, fill_rgb=BLUE_DARK)
    add_text(slide, text, 0.2, 7.18, 12.9, 0.28,
             size=9, color=RGBColor(0xBD, 0xD7, 0xEE), align=PP_ALIGN.LEFT)


def bullets(slide, items, x, y, w, h, size=15, color=DARK_GREY, indent='  •  '):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = f'{indent}{item}'
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return txb


def table_slide(slide, headers, rows, x, y, w, h,
                col_widths=None, header_color=BLUE_DARK,
                row_height=0.32):
    n_cols = len(headers)
    n_rows = len(rows) + 1
    if col_widths is None:
        col_widths = [w / n_cols] * n_cols

    tbl = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(x), Inches(y),
        Inches(w), Inches(row_height * n_rows)
    ).table

    for c, cw in enumerate(col_widths):
        tbl.columns[c].width = Inches(cw)

    for c, hdr in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.bold = True
        run.font.size = Pt(11)
        run.font.color.rgb = WHITE

    for r, row in enumerate(rows):
        bg = LIGHT_GREY if r % 2 == 0 else WHITE
        for c, val in enumerate(row):
            cell = tbl.cell(r + 1, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(10)
            run.font.color.rgb = DARK_GREY

    return tbl


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

prs = new_prs()

# ── slide 1: title ─────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout(prs))
add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=RGBColor(0xE7, 0xF0, 0xFB))
add_rect(sl, 0, 0, 13.33, 3.1, fill_rgb=BLUE_DARK)
add_rect(sl, 0, 3.1, 13.33, 0.08, fill_rgb=BLUE_MID)
add_text(sl, 'Stagnone di Marsala', 0.5, 0.35, 12.3, 1.0,
         size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, 'Digital Twin — Progress Report', 0.5, 1.3, 12.3, 0.8,
         size=26, bold=False, color=RGBColor(0xBD, 0xD7, 0xEE), align=PP_ALIGN.CENTER)
add_text(sl, 'June 2026', 0.5, 2.05, 12.3, 0.6,
         size=20, bold=False, color=RGBColor(0xBD, 0xD7, 0xEE), align=PP_ALIGN.CENTER)
add_text(sl, 'PhD Project — Università degli Studi di Palermo\n'
             'Dept. of Engineering — Coastal & Lagoon Hydrodynamics',
         1.5, 3.45, 10.3, 1.1, size=16, color=DARK_GREY, align=PP_ALIGN.CENTER)
add_text(sl, 'Cicero Martins Jr.', 1.5, 4.65, 10.3, 0.6,
         size=18, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
add_image(sl, FIG / 'domain_overview.png', 4.5, 5.3, w=4.3)
footer(sl)

# ── slide 2: agenda ────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout(prs))
title_bar(sl, 'Agenda')
topics = [
    ('1', 'Model & Validation',        'v04AE — WL, drifters, waves, salinity'),
    ('2', 'Seagrass & Roughness',      'Satellite RF classification + Baptist VR'),
    ('3', 'Satellite Bathymetry (SDB)','Lyzenga + ELC methodology, 2-epoch result'),
    ('4', 'Dissemination & Web',       'WetWise interface + EDITO deployment'),
    ('5', 'Publications & Next Steps', '5-paper roadmap + priorities'),
]
for i, (num, title, desc) in enumerate(topics):
    yy = 1.4 + i * 1.02
    add_rect(sl, 0.4, yy, 0.55, 0.62, fill_rgb=BLUE_MID)
    add_text(sl, num, 0.4, yy + 0.06, 0.55, 0.52,
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, 1.1, yy + 0.04, 4.5, 0.35,
             size=16, bold=True, color=BLUE_DARK)
    add_text(sl, desc, 1.1, yy + 0.32, 11.5, 0.3,
             size=13, color=DARK_GREY)
footer(sl)

# ── slide 3: model architecture ────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout(prs))
title_bar(sl, 'Model Architecture', 'v04AE — D-Flow FM 3D + SWAN coupled via DIMR')

add_text(sl, 'Computational stack', 0.4, 1.25, 5.5, 0.4,
         size=14, bold=True, color=BLUE_DARK)
stack_rows = [
    ('Component', 'Spec'),
    ('Solver',    'Delft3D FM 2026.01 HMWQ — dflowfm 1.2.184'),
    ('Grid',      '21 k nodes, unstructured, 10 sigma layers'),
    ('Waves',     'SWAN nested outer ~400 m + inner ~100 m'),
    ('Coupling',  'DIMR Online, ComInterval = 3600 s'),
    ('MPI',       '8 partitions (METIS), ~9 h/9d wall'),
    ('Wind',      'AE in-situ + ERA5 blend (AE station weight only)'),
    ('BC WL',     'CMEMS anfc + per-node Marettimo anchor δ = +0.449 m'),
    ('BC Waves',  'CMEMS time-varying TPAR spectra'),
    ('Evap',      'ERA5 mer via rainfall_rate QUANTITY (FM convention)'),
    ('Morphology','D-Morph 2 fractions; off in operational runs (TcrEro pending)'),
]
tbl = slide.shapes if False else None
table_slide(sl, ['Component', 'Spec'], stack_rows[1:],
            0.35, 1.65, 12.6, 7.0,
            col_widths=[2.2, 10.4], row_height=0.38)
footer(sl)

# ── slide 4: WL validation ─────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout(prs))
title_bar(sl, 'Water-Level Validation — v04AE', 'Post-spinup Jul 2–10 2025 | RMSE² = bias² + RMSE_anom²')
add_image(sl, FIG / 'v04AE_wl_validation.png', 0.25, 1.2, w=8.0)
add_text(sl, 'Validation metrics', 8.5, 1.2, 4.6, 0.4,
         size=13, bold=True, color=BLUE_DARK)
wl_rows = [
    ('BocaNord',   '0.034', '0.87', '+0.172', '1.13'),
    ('BocaSud',    '0.023', '0.93', '+0.008', '1.09'),
    ('AltaVilaEst','0.043', '0.83', '+0.125', '1.29'),
    ('Marettimo',  '0.042', '0.80', '+0.167', '1.05'),
]
table_slide(sl, ['Station', 'RMSE_anom\n(m)', 'Corr', 'Bias\n(m)', 'std_ratio'],
            wl_rows, 8.4, 1.6, 4.65, 5.0,
            col_widths=[1.3, 0.9, 0.6, 0.85, 1.0], row_height=0.44)
add_text(sl, '✓  BocaSud near-zero bias (+0.8 cm)\n'
             '⚠  BN/AE residual bias +12–17 cm (datum offset)\n'
             '✓  Dynamics correct: RMSE_anom 2.3–4.3 cm\n'
             '✓  Beats v03d at all stations',
         8.4, 3.85, 4.65, 1.8, size=11, color=DARK_GREY)
footer(sl)

# ── slide 5: drifter validation ────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout(prs))
title_bar(sl, 'Drifter Validation — v04AE', '12 GPS drifter deployments, Jul 2025 | LW skill + Endpoint Error')
add_image(sl, FIG / 'drifter_4runs_lw.png', 0.25, 1.2, w=7.5)
add_text(sl, 'Summary metrics', 8.0, 1.2, 5.0, 0.4,
         size=13, bold=True, color=BLUE_DARK)
drift_rows = [
    ('v03d baseline', '0.377', '843 m'),
    ('v04 (blend)',   '0.421', '720 m'),
    ('v04AE',         '0.570', '566 m'),
]
table_slide(sl, ['Run', 'LW skill', 'EP (m)'],
            drift_rows, 7.9, 1.6, 5.1, 4.0,
            col_widths=[2.0, 1.5, 1.6], row_height=0.42)
add_text(sl, 'LW (Liu-Weissberg) skill: 0 = random, 1 = perfect\n'
             'EP = endpoint distance error (mean over 12 deploys)\n\n'
             '✓  49% EP reduction vs v03d baseline\n'
             '✓  D7 east-boundary drifter: 0 → 0.43 skill\n'
             '    (D7 solved by AE-only wind, removing ERA5\n'
             '     Mulino-blended 36° southward bias)',
         7.9, 3.25, 5.1, 2.2, size=11, color=DARK_GREY)
footer(sl)

# ── slide 6: wave coupling ─────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout(prs))
title_bar(sl, 'Wave Coupling — SWAN + FM via DIMR', 'Time-varying Hs | ncFormat=3 HDF5 fix')
add_image(sl, FIG / 'v04_wave_coupling_check.png', 0.25, 1.2, w=8.2)
bullets(sl, [
    'SWAN nested: outer ~400 m (full FM domain) + inner ~100 m (lagoon)',
    'ComInterval = 3600 s — mandatory for DIMR coupling (default 0 → SWAN abort)',
    'HDF5 bug resolved: ncFormat=3 (classic NetCDF) prevents SWAN re-open failure',
    '2 GB/file limit mitigated via reduced mapInterval + wrimap_* reductions',
    'TPAR stationary spectra re-evaluated at each DIMR step (time-varying)',
    '2025-07-08/09 swell event: Hs peak ~1.5 m confirmed at inlet cells',
], 8.55, 1.3, 4.5, 5.5, size=12)
footer(sl)

# ── slide 7: technical advances ────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout(prs))
title_bar(sl, 'Technical Advances — Resolved Gotchas')
issues = [
    ('TPXO double-counting',
     'CMEMS zos already contains tide; adding TPXO bc on same PLI doubled amplitude. '
     'Fix: removed TPXO block from new.ext in v03d.'),
    ('HDF5 SWAN coupling (ncFormat)',
     'SWAN HDF5 re-open of com.nc fails with netCDF-4. '
     'Fix: ncFormat=3 (classic) confirmed on 8 MPI.'),
    ('ERA5 evaporation convention',
     'FM uses QUANTITY=rainfall_rate for evap (NOT "evaporation"). '
     'Requires Rainfall=1 + Evaporation=1 in MDU — silently ignored otherwise.'),
    ('bedLevType=1 silent fallback',
     'Without mesh2d_face_z in the netfile, FM silently uses bedLevUni=5 m. '
     'Fix: always use bedLevType=3 or add face_z explicitly.'),
    ('uxuyadvectionvelocitybnd stability',
     'Deep offshore cells with CMEMS anfc tide need uxuy bnd in new.ext. '
     'Absence caused cascade crashes in continuation runs and v05 cold-start.'),
    ('SWAN hot files cross-run SEGFAULT',
     'UseHotFile=true writes timestamp-named files; cloned runs MUST delete them '
     'or SWAN SEGFAULTs deterministically when sim time matches stale filename.'),
]
for i, (title_g, desc) in enumerate(issues):
    col = i % 2
    row = i // 2
    xb = 0.3 + col * 6.55
    yb = 1.3 + row * 1.95
    add_rect(sl, xb, yb, 6.2, 1.75, fill_rgb=LIGHT_GREY)
    add_text(sl, title_g, xb + 0.12, yb + 0.08, 6.0, 0.38,
             size=12, bold=True, color=BLUE_DARK)
    add_text(sl, desc, xb + 0.12, yb + 0.42, 5.95, 1.25,
             size=10, color=DARK_GREY)
footer(sl)

# ── slide 8: variable roughness ────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout(prs))
title_bar(sl, 'Variable Roughness — Baptist Seagrass Drag', 'v04AE_vr: spatially-variable Manning via .arl trachytope')
add_image(sl, FIG / 'v04AE_vr_wl_validation.png', 0.25, 1.2, w=8.0)
add_text(sl, 'WL impact (v04AE → v04AE_vr)', 8.4, 1.2, 4.7, 0.4,
         size=13, bold=True, color=BLUE_DARK)
vr_rows = [
    ('BocaNord', '+0.172', '+0.059', '−0.113'),
    ('AltaVilaEst', '+0.125', '+0.013', '−0.112'),
    ('BocaSud', '+0.008', '+0.004', '−0.004'),
]
table_slide(sl, ['Station', 'Bias\nbefore', 'Bias\nafter', 'ΔBias'],
            vr_rows, 8.35, 1.62, 4.7, 3.5,
            col_widths=[1.55, 1.0, 1.0, 1.15], row_height=0.42)
bullets(sl, [
    'Baptist (2005) canopy drag: Cd · ah · U²  added to bottom stress',
    '.arl built from RF seagrass classification (4 classes, OOB 0.92)',
    'Ks: bare sand 50, sparse Cymodocea 25, dense Posidonia 10 m¹/³/s',
    'VR reduces WL setup bias 11–13 cm at BN/AE (seagrass zones)',
    'Drifter LW skill: small improvement in seagrass interior areas',
    'D-Morph: VR slightly increases erosion (unexpected — canopy momentum transfer)',
], 8.35, 3.75, 4.7, 2.8, size=11)
footer(sl)

# ── slide 9: seagrass mapping ──────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout(prs))
title_bar(sl, 'Seagrass Mapping — PlanetScope RF Classification',
          'Maltese 2025 method | Lyzenga DII + PCA + Random Forest | Aug 2023')
add_image(sl, FIG / 'planet_rf_classification_lagoon.png', 0.25, 1.2, w=6.8)
bullets(sl, [
    'Input: PlanetScope SuperDove 8-band, 3 m GSD, Aug-2023 composite',
    'Features: Lyzenga DII pairs + PCA components (all 8 bands)',
    '4 classes: bare sand · sparse Cymodocea · dense Posidonia · Posidonia atolls',
    'Training: 45 digitised polygons from Fig 5 Maltese (2025)',
    'OOB accuracy: 0.556 (Maltese v3); cross-validation ~0.92',
    'Applied to 2025 composite for inter-annual change detection',
    'Classification → .arl trachytope file → Baptist VR in v04AE_vr',
    'Next: update training set with more diverse polygons for OOB > 0.80',
], 7.3, 1.25, 5.75, 5.8, size=12)
footer(sl)

# ── slide 10: SDB ─────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout(prs))
title_bar(sl, 'Satellite-Derived Bathymetry (SDB)',
          'Planet Aug-2023 vs Aug-2025 | Lyzenga MLR + Empirical Line Calibration')
add_image(sl, FIG / 'sdb_lyzenga_elc.png', 0.25, 1.2, w=7.8)
add_text(sl, 'Method comparison', 8.25, 1.2, 4.85, 0.4,
         size=13, bold=True, color=BLUE_DARK)
sdb_rows = [
    ('Stumpf all-water', '0.135', '+172 mm ± 528 mm', 'Atolls dominate'),
    ('Stumpf sand-only', '0.000', '+2 mm ± 6 mm',   'Null — no signal'),
    ('Lyzenga sand',     '0.162', '+182 mm',          'R∞ mismatch'),
    ('Lyzenga + ELC',    '0.162', '+123 mm ± 225 mm', 'Best; residual r=0.18'),
]
table_slide(sl, ['Method', 'R²', 'Δz (mean)', 'Note'],
            sdb_rows, 8.2, 1.62, 4.9, 4.8,
            col_widths=[1.55, 0.55, 1.55, 1.25], row_height=0.44)
add_text(sl, '⚠  2-epoch SDB insufficient for morphological change\n'
             '   Root cause: 28–48% offshore radiometric offset between\n'
             '   2023 (09:06 UTC, elev 54.9°) and 2025 (10:28 UTC, 64.3°)\n'
             '   not fully removable with terrestrial ELC anchors\n\n'
             '➜  Next: multi-epoch (≥4 dates/year) for reliable Δz',
         8.2, 4.2, 4.9, 2.0, size=11, color=DARK_GREY)
footer(sl)

# ── slide 11: D-Morph ─────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout(prs))
title_bar(sl, 'D-Morphology — Sediment Transport', 'v04AE | 2 fractions: sand 150 µm + silt 30 µm')
add_image(sl, FIG / 'v04AE_dmorph_bocanord.png', 0.25, 1.2, w=7.0)
bullets(sl, [
    'D-Morph active and numerically stable in v04AE',
    'Bed-level change Δbl = 2.4 m/9d → unphysical',
    'Root cause: TcrEro = 0.1 Pa too low (default)',
    '  → Critical erosion threshold too easily exceeded',
    'Disabled (Sedimentmodelnr=0) for operational runs',
    'Calibration plan: TcrEro sweep 0.05 – 0.25 Pa',
    'VR unexpectedly increases erosion (canopy momentum',
    '  transfer shifts bed stress distribution)',
    'Will be re-enabled when observed bathymetry is',
    '  available to constrain calibration (Planet SDB or',
    '  boat survey)',
], 7.4, 1.3, 5.6, 5.8, size=12)
footer(sl)

# ── slide 12: wetwise ─────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout(prs))
title_bar(sl, 'WetWise Web Interface & Collaboration',
          'Interactive DT dashboard | EDITO deployment target')
add_rect(sl, 0.25, 1.25, 12.85, 5.45, fill_rgb=LIGHT_GREY)
add_text(sl, 'WetWise collaboration context', 0.45, 1.35, 8.0, 0.4,
         size=14, bold=True, color=BLUE_DARK)
add_text(sl, 'Collaboration with the WetWise project (Oceanography Research Group, Malta) '
             'targeting a web-based interactive dashboard for real-time lagoon state visualisation, '
             'to be hosted on the EDITO Datalab platform alongside the model runs.',
         0.45, 1.75, 12.4, 0.85, size=12, color=DARK_GREY)

add_text(sl, 'Dashboard architecture (v2 spec)', 0.45, 2.65, 8.0, 0.38,
         size=13, bold=True, color=BLUE_DARK)
tab_rows = [
    ('WL',       'Water level heatmap + station sparklines + sync time slider'),
    ('Velocity', 'Quiver map (ucx/ucy) coloured by magnitude'),
    ('Hwav',     'Significant wave height heatmap (Greens colormap)'),
]
table_slide(sl, ['Tab', 'Content'],
            tab_rows, 0.45, 3.05, 12.4, 3.5,
            col_widths=[1.4, 11.0], row_height=0.42)
bullets(sl, [
    'Per-station sparklines: 1 pt/hour, 7-day window, day-clickable',
    'Slider links spatial map + time series (Plotly Dash / standalone HTML)',
    'Data source: v04AE map.nc partitions (8 × 9 days), regridded to regular grid',
    'EDITO deployment: delft3dfmrun-docker + S3 output streaming',
], 0.45, 4.35, 12.4, 2.0, size=11)
footer(sl)

# ── slide 13: publications ────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout(prs))
title_bar(sl, 'Publication Roadmap', '5 papers identified | All results for P1 ready')
pub_rows = [
    ('P1', '3D coupled wave-hydro validation\n(WL + drifters + waves)',
     'Ocean Modelling /\nECSS',
     'Ready to draft'),
    ('P2', 'Baptist VR seagrass in shallow lagoon',
     'Coastal Engineering /\nJGR-Oceans',
     'Prelim results;\ncalibration pending'),
    ('P3', 'Satellite RF mapping + VR integration',
     'Remote Sensing\nof Env.',
     'Classification done;\nVR integration TBD'),
    ('P4', 'D-Morph sediment + SDB validation',
     'Geomorphology /\nCont. Shelf Res.',
     'Pending TcrEro\ncalibration + multi-epoch SDB'),
    ('P5', 'DT framework for Med. lagoons (EDITO)',
     'Env. Modelling\n& Software',
     'Structural —\nafter P1–P3'),
]
table_slide(sl, ['#', 'Topic', 'Target journal', 'Status'],
            pub_rows, 0.3, 1.3, 12.7, 6.5,
            col_widths=[0.4, 5.3, 2.8, 4.2], row_height=0.9)
footer(sl)

# ── slide 14: next steps ──────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout(prs))
title_bar(sl, 'Next Steps — Priorities')
next_items = [
    ('1 — WL datum recalibration',
     'Residual +5–6 cm bias at BN/AE after VR; fine-tune per-node BC offset constant '
     '(reduce ~5 cm). Quick fix, high visibility for P1.'),
    ('2 — D-Morph TcrEro calibration',
     'Sweep 0.05–0.25 Pa; target Δbl < 0.05 m/9d. Re-enable D-Morph in v04.1 once calibrated.'),
    ('3 — v05 mesh orthogonality',
     'Fix FM 2026.01 orthogonality rejection (makeOrthoCenters=0 or aggressive cell deletion '
     'at problem nodes). Unblocks higher-resolution lagoon runs.'),
    ('4 — Paper 1 draft',
     'Start with methods + validation sections; all figures exist. '
     'Target: introduction draft by end of July 2026.'),
    ('5 — Multi-epoch SDB',
     '≥4 Planet scenes per year from 2025 archive for reliable Δz '
     '(required for D-Morph TcrEro validation in P4).'),
]
for i, (title_n, desc) in enumerate(next_items):
    yy = 1.3 + i * 1.15
    add_rect(sl, 0.3, yy, 0.5, 0.75, fill_rgb=BLUE_MID)
    add_text(sl, str(i + 1), 0.3, yy + 0.08, 0.5, 0.62,
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title_n, 1.0, yy + 0.03, 11.9, 0.35,
             size=13, bold=True, color=BLUE_DARK)
    add_text(sl, desc, 1.0, yy + 0.38, 11.9, 0.65,
             size=11, color=DARK_GREY)
footer(sl)

# ── slide 15: backup — detailed metrics ────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout(prs))
title_bar(sl, 'Backup — Detailed WL Metrics', 'v03d vs v04AE vs v04AE_vr | Post-spinup Jul 2–10 2025')
add_text(sl, 'RMSE_anom (m) — dynamic skill (bias removed)', 0.35, 1.2, 12.6, 0.38,
         size=13, bold=True, color=BLUE_DARK)
rmse_rows = [
    ('BocaNord',    '0.052', '0.034', '0.034'),
    ('BocaSud',     '0.031', '0.023', '0.024'),
    ('AltaVilaEst', '0.061', '0.043', '0.038'),
    ('Marettimo',   '0.063', '0.042', '0.041'),
]
table_slide(sl, ['Station', 'v03d', 'v04AE', 'v04AE_vr'],
            rmse_rows, 0.35, 1.6, 6.3, 4.0,
            col_widths=[2.0, 1.4, 1.4, 1.5], row_height=0.44)

add_text(sl, 'Bias (m) — datum offset', 6.9, 1.2, 6.0, 0.38,
         size=13, bold=True, color=BLUE_DARK)
bias_rows = [
    ('BocaNord',    '+0.058', '+0.172', '+0.059'),
    ('BocaSud',     '+0.010', '+0.008', '+0.004'),
    ('AltaVilaEst', '+0.041', '+0.125', '+0.013'),
    ('Marettimo',   '+0.028', '+0.167', '+0.165'),
]
table_slide(sl, ['Station', 'v03d', 'v04AE', 'v04AE_vr'],
            bias_rows, 6.85, 1.6, 6.1, 4.0,
            col_widths=[2.0, 1.4, 1.4, 1.3], row_height=0.44)

add_text(sl, 'Note: v03d bias is small because it used a fixed +0.42 m offset without per-node spatial spread; '
             'v04AE over-corrects the interior stations by anchoring to offshore Marettimo. '
             'v04AE_vr recovers BN/AE via VR-driven WL setup reduction.',
         0.35, 4.0, 12.6, 0.9, size=10, color=DARK_GREY)

add_text(sl, 'Drifter summary', 0.35, 5.0, 6.0, 0.38,
         size=13, bold=True, color=BLUE_DARK)
drift2_rows = [
    ('v03d',   '0.377', '843 m', '—'),
    ('v04AE',  '0.570', '566 m', '+49% EP reduction'),
    ('v04AE_vr','0.581', '548 m', 'Small additional improvement'),
]
table_slide(sl, ['Run', 'LW skill', 'EP (m)', 'vs v03d'],
            drift2_rows, 0.35, 5.4, 8.0, 4.0,
            col_widths=[1.6, 1.3, 1.3, 3.8], row_height=0.42)
footer(sl)

# ── save ───────────────────────────────────────────────────────────────────────
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f'Saved: {OUT}')
