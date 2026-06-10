"""
Generate Wetwise portal presentation (6 slides) as PowerPoint.
Output: outputs/wetwise_stagnone_panel_proposal.pptx

Requires: python-pptx  (pip install python-pptx)
"""
from __future__ import annotations
import io
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as pu

OUT = Path(__file__).parent.parent / 'outputs' / 'wetwise_stagnone_panel_proposal.pptx'
OUT.parent.mkdir(parents=True, exist_ok=True)

# ── Colour palette ─────────────────────────────────────────────────────────────
def rgb(h): return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))
DARK   = rgb('0d2137')
MID    = rgb('1a6b8a')
TEAL   = rgb('2ba3a8')
SAND   = rgb('e8d5a3')
LIGHT  = rgb('f0f6fa')
WHITE  = rgb('ffffff')
GREY   = rgb('8a9bb0')
GREEN  = rgb('3aaa6e')
ORANGE = rgb('e07b3a')
RED    = rgb('c0392b')
BGSLIDE= rgb('f0f6fa')

SW, SH = Inches(16), Inches(9)   # widescreen 16×9

# ── pptx helpers ───────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h))   # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, size=18, bold=False, color=None,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    color = color or DARK
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic    = italic
    return tb


def header(slide, title, subtitle=None):
    """Standard dark header bar + optional subtitle."""
    add_rect(slide, 0, 0, 16, 1.1, DARK)
    add_text(slide, title, 0.35, 0.10, 13, 0.6,
             size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.65, 13, 0.38,
                 size=14, color=TEAL)
    # footer
    add_rect(slide, 0, 8.65, 16, 0.35, DARK)
    add_text(slide, 'Stagnone di Marsala Digital Twin  ·  Università degli Studi di Palermo',
             0.2, 8.67, 11, 0.28, size=10, color=GREY)
    add_text(slide, 'June 2026', 14.5, 8.67, 1.3, 0.28,
             size=10, color=GREY, align=PP_ALIGN.RIGHT)


def mpl_to_slide(slide, fig, x, y, w, h):
    """Render a matplotlib figure and embed as image in slide."""
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor=fig.get_facecolor())
    buf.seek(0)
    slide.shapes.add_picture(buf, Inches(x), Inches(y), Inches(w), Inches(h))
    plt.close(fig)


def bullet_box(slide, items, x, y, w, h, size=13, color=None,
               spacing_pt=22, indent='▸  '):
    color = color or DARK
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(spacing_pt if i > 0 else 0)
        run = p.add_run()
        run.text = indent + item
        run.font.size  = Pt(size)
        run.font.color.rgb = color


# ═══════════════════════════════════════════════════════════════════════════════
def slide_1_title(prs):
    slide = blank_slide(prs)
    bg(slide, DARK)
    # accent side panel
    add_rect(slide, 0, 0, 5.5, 9, MID, line_color=None)
    add_rect(slide, 0, 0, 5.5, 9, rgb('000000'))   # reuse as semi-transparent via opacity trick
    # actually just do a slightly lighter blue strip
    s = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(5.2), Inches(9))
    s.fill.solid(); s.fill.fore_color.rgb = MID
    s.fill.fore_color.theme_color  # no-op
    s.line.fill.background()
    # overlay darker tint
    s2 = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(5.2), Inches(9))
    s2.fill.solid(); s2.fill.fore_color.rgb = DARK
    from pptx.util import Pt as PT
    s2.line.fill.background()
    # (pptx doesn't do opacity easily; just use a slightly lighter blue strip)
    s.fill.fore_color.rgb = rgb('163d5c')

    # accent vertical bar
    add_rect(slide, 5.55, 2.8, 0.07, 2.2, TEAL)

    add_text(slide, 'Stagnone di Marsala',  5.75, 2.0, 9.5, 1.0,
             size=38, bold=True, color=WHITE)
    add_text(slide, 'Digital Twin — 3D Coupled Hydrodynamic Model',
             5.75, 3.05, 9.5, 0.7, size=22, color=SAND)
    add_text(slide, 'Proposal: Hydrodynamic Panel for the Wetwise Portal',
             5.75, 3.80, 9.5, 0.6, size=18, color=TEAL)
    add_text(slide, 'Cicero Martins Jr.  ·  Università degli Studi di Palermo  ·  June 2026',
             5.75, 5.20, 9.5, 0.5, size=13, color=GREY)

    # domain schematic on the left panel
    fig, ax = plt.subplots(figsize=(4, 7), facecolor='#0d2137')
    ax.set_facecolor('#0d2137'); ax.set_xlim(0,1); ax.set_ylim(0,1); ax.axis('off')
    ax.add_patch(plt.Rectangle((0,0),1,1, color='#1a4a6a'))
    coast_x = [0.50,0.55,0.62,0.65,0.68,0.70,0.71,0.69,0.67,0.63,0.58,0.52,0.46,0.42,0.40,0.41,0.44,0.48,0.50]
    coast_y = [0.97,0.93,0.90,0.85,0.78,0.68,0.58,0.48,0.38,0.28,0.20,0.14,0.13,0.17,0.27,0.38,0.50,0.64,0.78]
    ax.fill(coast_x + [1,1,0,0], coast_y + [0,1,1,0.97], color='#2a5a2a', zorder=1)
    ax.fill([0.20,0.28,0.30,0.32,0.33,0.32,0.29,0.25,0.21,0.19,0.20],
            [0.88,0.87,0.80,0.70,0.57,0.42,0.28,0.20,0.22,0.38,0.62],
            color='#2ba3a8', alpha=0.75, zorder=2)
    ax.text(0.25, 0.52, 'Stagnone\nLagoon', color='white', fontsize=11,
            ha='center', fontweight='bold', zorder=3)
    ax.add_patch(plt.Circle((0.10, 0.42), 0.06, color='#2a5a2a', zorder=2))
    ax.text(0.10, 0.42, 'Marettimo', color='#ccc', fontsize=7, ha='center', zorder=3)
    ax.text(0.5, 0.08, 'Sicily — W Coast', color='#aaddaa', fontsize=8, ha='center')
    mpl_to_slide(slide, fig, 0.2, 0.8, 4.8, 7.4)
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
def slide_2_model(prs):
    slide = blank_slide(prs)
    bg(slide, BGSLIDE)
    header(slide, 'The Digital Twin Model',
           'Delft3D FM 3D + SWAN coupled · validated July 2025')

    # LEFT: domain schematic
    fig, ax = plt.subplots(figsize=(5.5, 6.5), facecolor='#ddeeff')
    ax.set_facecolor('#b8d8f0'); ax.set_xlim(0,1); ax.set_ylim(0,1); ax.axis('off')
    ax.add_patch(plt.Rectangle((0,0),1,1, color='#c0daf0'))
    cst_x = [0.50,0.55,0.60,0.63,0.65,0.66,0.67,0.66,0.64,0.61,0.57,0.52,0.47,0.43,0.42,0.43,0.46,0.50]
    cst_y = [0.97,0.94,0.90,0.84,0.76,0.65,0.54,0.43,0.32,0.22,0.15,0.10,0.10,0.15,0.26,0.40,0.57,0.75]
    ax.fill(cst_x+[1,1,0,0], cst_y+[0,1,1,0.97], color='#c8d8a8', zorder=1)
    lag_x = [0.25,0.30,0.33,0.35,0.37,0.38,0.37,0.35,0.32,0.28,0.25,0.23,0.23,0.24,0.25]
    lag_y = [0.88,0.86,0.81,0.73,0.62,0.50,0.37,0.26,0.18,0.14,0.16,0.27,0.45,0.65,0.80]
    ax.fill(lag_x, lag_y, color='#4ab0c8', alpha=0.85, zorder=2)
    ax.add_patch(plt.Circle((0.11, 0.55), 0.055, color='#c8d8a8', zorder=2))
    ax.text(0.11, 0.55, 'Marettimo\n(BC anchor)', color='#333', fontsize=7.5, ha='center', zorder=3)
    ax.text(0.30, 0.50, 'Stagnone\n11×3 km', color='white', fontsize=9, ha='center', fontweight='bold', zorder=3)
    for bx, by, lbl in [(0.27,0.87,'BN'),(0.25,0.19,'BS'),(0.36,0.50,'AE')]:
        ax.plot(bx, by, 'v', color='#e07b3a', ms=10, zorder=5)
        ax.text(bx+0.04, by, lbl, color='#e07b3a', fontsize=8, zorder=5)
    ax.annotate('15 m mesh\n(lagoon)', xy=(0.35,0.62), xytext=(0.65,0.74),
                arrowprops=dict(arrowstyle='->', color='#333', lw=1.2),
                fontsize=8, color='#333', ha='center')
    ax.annotate('~300 m\n(offshore)', xy=(0.10,0.30), xytext=(0.02,0.17),
                arrowprops=dict(arrowstyle='->', color='#333', lw=1.2),
                fontsize=8, color='#333')
    ax.set_title('Model domain  —  UTM33N / WGS84', fontsize=10, pad=3, color='#333')
    mpl_to_slide(slide, fig, 0.3, 1.2, 6.2, 7.2)

    # RIGHT: specs table
    specs = [
        ('Kernel',       'Delft3D FM 3D  +  SWAN (DIMR Online coupled)'),
        ('Hydrodynamics','D-Flow FM 2026.01 HMWQ  ·  kernel 1.2.184'),
        ('Waves',        'SWAN 40.91AB  ·  outer ~400 m / inner ~100 m'),
        ('Domain',       'Lagoon 11×3 km + offshore + Marettimo'),
        ('Resolution',   '15 m (lagoon interior)  →  300 m (offshore)'),
        ('Layers',       '10 sigma layers (3D)  ·  hypersaline IC ~42 ppt'),
        ('Forcings',     'CMEMS Med reanalysis/anfc + ERA5 + local wind'),
        ('BC offset',    'Per-node WL offset anchored to Marettimo TAD'),
        ('Roughness',    'Satellite seagrass map → Baptist trachytopes'),
        ('Validated',    'July 2025 · 9-day · WL corr BN=0.98, BS=0.90'),
    ]
    add_rect(slide, 6.8, 1.2, 8.8, 7.2, WHITE, line_color=GREY, line_width=Pt(0.5))
    for i, (k, v) in enumerate(specs):
        y = 1.45 + i * 0.67
        add_text(slide, k + ':', 7.0, y, 2.4, 0.55, size=11, bold=True, color=MID)
        add_text(slide, v,       9.5, y, 5.9, 0.55, size=11, color=DARK)
        if i < len(specs)-1:
            add_rect(slide, 6.85, y+0.52, 8.7, 0.02, LIGHT)
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
def slide_3_outputs(prs):
    slide = blank_slide(prs)
    bg(slide, BGSLIDE)
    header(slide, 'Available Outputs & Validation',
           'Variables from July 2025 validated run  ·  RMSE² = bias² + RMSE_anom²')

    # Variable cards
    vars_ = [
        ('Water Level (WL)', '1a6b8a', 'his.nc + map.nc\n10-min resolution\nMSL-referenced'),
        ('Currents (u, v)',  '2ba3a8', 'map.nc  3D + 2D\n10-min resolution\nvector fields'),
        ('Wave Height Hwav', '1d7a7e', 'SWAN → com.nc\nStationary / 3600 s\nTp, Dir available'),
        ('Salinity (sa1)',   '1a5060', 'map.nc  3D\nhypersaline IC\n~42 ppt interior'),
    ]
    for i, (name, col_h, desc) in enumerate(vars_):
        x = 0.35 + i * 3.88
        add_rect(slide, x, 1.25, 3.5, 2.0, rgb(col_h))
        add_text(slide, name, x+0.12, 1.30, 3.2, 0.55, size=13, bold=True, color=WHITE)
        add_text(slide, desc, x+0.12, 1.90, 3.2, 1.20, size=11, color=rgb('d0eef5'))

    # Validation table
    add_text(slide, 'Validation — July 2025 (post-spinup window: day 1 → 9)',
             0.35, 3.45, 12, 0.55, size=14, bold=True, color=DARK)

    headers = ['Station', 'Variable', 'Bias (m)', 'RMSE (m)', 'RMSE_anom (m)', 'Corr (r)', 'Status']
    rows = [
        ['BocaNord',    'WL', '−0.01', '0.041', '0.041', '0.98', '✔ excellent'],
        ['BocaSud',     'WL', '+0.03', '0.064', '0.057', '0.90', '✔ good'],
        ['AltaVilaEst', 'WL', '+0.09', '0.162', '0.133', '0.36', '⚠ surge-dominated'],
        ['Marettimo',   'WL', '+0.06', '0.089', '0.065', '0.95', '✔ good (offshore)'],
    ]
    col_x = [0.35, 2.45, 4.30, 6.00, 7.80, 9.80, 11.50]
    col_w = [2.0,  1.8,  1.6,  1.7,  1.9,  1.5,  4.0]

    add_rect(slide, 0.35, 4.00, 15.30, 0.55, MID)   # header row
    for j, (h, cw) in enumerate(zip(headers, col_w)):
        add_text(slide, h, col_x[j]+0.05, 4.03, cw, 0.48,
                 size=11, bold=True, color=WHITE)

    row_bg = [WHITE, LIGHT]
    for i, row in enumerate(rows):
        y = 4.60 + i * 0.72
        add_rect(slide, 0.35, y, 15.30, 0.65, row_bg[i % 2],
                 line_color=LIGHT, line_width=Pt(0.3))
        for j, (val, cw) in enumerate(zip(row, col_w)):
            txt_color = GREEN if (j == 5 and float(val) > 0.8) else (
                        ORANGE if j == 6 and '⚠' in val else DARK)
            add_text(slide, val, col_x[j]+0.05, y+0.08, cw, 0.50,
                     size=11, color=txt_color)

    add_text(slide,
             '⚠  AltaVilaEst: shallow interior station, WL dominated by wind setup / surge — '
             'tidal signal weak; skill will improve with variable roughness calibration.',
             0.35, 8.55, 15.3, 0.38, size=9.5, color=ORANGE, italic=True)
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
def slide_4_panel(prs):
    slide = blank_slide(prs)
    bg(slide, BGSLIDE)
    header(slide, 'Proposed Panel Design',
           '3 interactive tabs · linked spatial map + station time series + sparkline navigator')

    # Tab bar
    tabs = [('Water Level', MID), ('Currents', GREY), ('Wave Height', GREY)]
    for i, (t, tc) in enumerate(tabs):
        add_rect(slide, 0.35 + i*2.15, 1.20, 2.0, 0.50, tc)
        add_text(slide, t, 0.40+i*2.15, 1.22, 1.9, 0.46, size=12, bold=(i==0),
                 color=WHITE, align=PP_ALIGN.CENTER)

    # Map mockup
    fig, axm = plt.subplots(figsize=(7, 5.5), facecolor='#ddeeff')
    axm.set_facecolor('#c0daf0'); axm.axis('off'); axm.set_xlim(0,1); axm.set_ylim(0,1)
    xx, yy = np.meshgrid(np.linspace(0,1,50), np.linspace(0,1,80))
    np.random.seed(42)
    wl  = np.sin(xx*3)*np.cos(yy*2)*0.2 + 0.03*np.random.randn(80,50)
    cm  = axm.pcolormesh(xx, yy, wl, cmap='RdBu_r', vmin=-0.25, vmax=0.25, rasterized=True)
    lag_x = [0.25,0.30,0.33,0.35,0.37,0.38,0.37,0.35,0.32,0.28,0.25,0.23,0.23,0.24,0.25]
    lag_y = [0.88,0.86,0.81,0.73,0.62,0.50,0.37,0.26,0.18,0.14,0.16,0.27,0.45,0.65,0.80]
    for bx, by, lbl in [(0.27,0.87,'BocaNord'),(0.25,0.19,'BocaSud'),(0.36,0.50,'AltaVilaEst')]:
        axm.plot(bx, by, 'v', color='#e07b3a', ms=10, zorder=5)
        axm.text(bx+0.03, by, lbl, color='#e07b3a', fontsize=8, va='center', zorder=5)
    plt.colorbar(cm, ax=axm, label='WL (m)', shrink=0.65, pad=0.01)
    axm.set_title('Spatial map  ·  Plotly heatmap (zoom / pan / hover)',
                  fontsize=10, color='#333', pad=3)
    axm.set_xlabel('← time slider →', fontsize=9, color='#555')
    mpl_to_slide(slide, fig, 0.35, 1.75, 7.8, 6.50)

    # Time series
    fig2, ax2 = plt.subplots(figsize=(6.5, 2.8), facecolor='#f0f6fa')
    t_arr = np.linspace(0,9,864); np.random.seed(7)
    wl_ts = 0.15*np.sin(t_arr*2*np.pi/0.52) + 0.07*np.sin(t_arr*2*np.pi/3) + 0.015*np.random.randn(864)
    ax2.plot(t_arr, wl_ts, color='#1a6b8a', lw=1.3)
    ax2.axvline(4.5, color='#e07b3a', lw=1.8, ls='--', label='selected time')
    ax2.set_xlabel('Day (Jul 2025)', fontsize=9); ax2.set_ylabel('WL (m)', fontsize=9)
    ax2.set_title('Station time series · 10-min resolution · click to select', fontsize=9)
    ax2.legend(fontsize=8); ax2.grid(alpha=0.25); ax2.set_facecolor('#f0f6fa')
    fig2.tight_layout()
    mpl_to_slide(slide, fig2, 8.40, 1.90, 7.30, 3.40)

    # Sparkline
    fig3, ax3 = plt.subplots(figsize=(6.5, 1.8), facecolor='#f0f6fa')
    t_hr = np.linspace(0,9,9*24); sp = 0.15*np.sin(t_hr*2*np.pi/0.52)+0.07*np.sin(t_hr*2*np.pi/3)
    for d in range(9):
        seg = (t_hr >= d) & (t_hr < d+1)
        ax3.plot(t_hr[seg], sp[seg], color='#1a6b8a' if d!=4 else '#e07b3a', lw=1.6)
    ax3.set_title('Sparkline — 7-day overview · 1 pt/hr · click any day to zoom', fontsize=9)
    ax3.tick_params(labelbottom=False); ax3.grid(alpha=0.2); ax3.set_facecolor('#f0f6fa')
    fig3.tight_layout()
    mpl_to_slide(slide, fig3, 8.40, 5.50, 7.30, 2.55)

    return slide


# ═══════════════════════════════════════════════════════════════════════════════
def slide_5_demo(prs):
    slide = blank_slide(prs)
    bg(slide, BGSLIDE)
    header(slide, 'Current State — July 2025 Demo',
           'V1 HTML demo operational · V2 Plotly bundle in progress')

    # Metric badges
    metrics = [
        ('r = 0.98', 'BocaNord', '3aaa6e'),
        ('r = 0.90', 'BocaSud',  '3aaa6e'),
        ('r = 0.95', 'Marettimo','2ba3a8'),
        ('skill 0.57','Drifters', '1a6b8a'),
    ]
    for i, (val, lbl, col) in enumerate(metrics):
        bx = 9.80 + i * 1.55
        add_rect(slide, bx, 1.25, 1.40, 1.15, rgb(col))
        add_text(slide, val, bx+0.05, 1.30, 1.30, 0.55,
                 size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, lbl, bx+0.05, 1.82, 1.30, 0.45,
                 size=10, color=WHITE, align=PP_ALIGN.CENTER)

    # Done column
    add_rect(slide, 0.35, 1.25, 7.0, 3.50, WHITE, line_color=GREEN, line_width=Pt(1.5))
    add_text(slide, '✔  What exists today', 0.50, 1.30, 6.7, 0.55,
             size=13, bold=True, color=GREEN)
    done = [
        '9-day coupled FM+SWAN run, Jul 7–16 2025',
        'WL at BocaNord / BocaSud / AltaVilaEst + Marettimo',
        'V1 HTML demo: 6 pre-rendered MP4 videos\n'
        '   (WL, ucx, ucy, ucmag, Hwav, salinity)',
        'Drifter validation: Lagrangian skill = 0.570',
        'Seagrass Baptist trachytopes integrated (v04AE_vr)',
        'Continuation chain validated (Jul 7→20 window)',
    ]
    bullet_box(slide, done, 0.50, 1.95, 6.7, 2.65,
               size=11.5, color=DARK, spacing_pt=16)

    # In progress column
    add_rect(slide, 7.70, 1.25, 7.95, 3.50, WHITE, line_color=MID, line_width=Pt(1.5))
    add_text(slide, '⚙  V2 improvements in progress', 7.85, 1.30, 7.65, 0.55,
             size=13, bold=True, color=MID)
    wip = [
        'Replace MP4 → Plotly interactive maps (zoom/pan/hover)',
        'Linked slider: spatial map + WL chart advance together',
        'Sparkline navigator: 1 pt/hr, click any day to zoom in',
        'Variable roughness (seagrass) validation vs baseline',
    ]
    bullet_box(slide, wip, 7.85, 1.95, 7.65, 2.65,
               size=11.5, color=DARK, spacing_pt=16)

    # Known limitations
    add_rect(slide, 0.35, 4.95, 15.30, 2.20, rgb('fff8f0'), line_color=ORANGE, line_width=Pt(1.0))
    add_text(slide, '⚠  Known limitations / active work', 0.55, 5.00, 14.8, 0.55,
             size=13, bold=True, color=ORANGE)
    limits = [
        'AltaVilaEst WL skill low (r=0.36): tide signal weak at interior shallow station — '
        'dominated by wind setup & surge; VR (variable roughness) calibration is the next lever.',
        'Seagrass trachytopes: Baptist params not yet calibrated (TcrEro conservative); '
        'D-Morph disabled for Paper 1 (Δbl=2.4m/9d unphysical). Planned for v05 with observed bathy.',
        'Offshore Posidonia artifact in roughness ARL: RF classifier extrapolates outside training '
        'domain offshore. Lagoon polygon mask fix in backlog.',
    ]
    bullet_box(slide, limits, 0.55, 5.65, 15.0, 1.90,
               size=10.5, color=DARK, spacing_pt=13)

    # Output path note
    add_text(slide,
             'Output path:  outputs/wetwise_demo_jul25/index.html  '
             '·  manifest.json  →  consumed by portal dev',
             0.35, 7.35, 15.3, 0.42, size=10, color=GREY, italic=True)
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
def slide_6_roadmap(prs):
    slide = blank_slide(prs)
    bg(slide, BGSLIDE)
    header(slide, 'Integration & Roadmap',
           'Operational pipeline · data contract for portal · next milestones')

    # Pipeline flow diagram
    steps = [
        ('CMEMS anfc\nMed 1/24°',      '1a6b8a'),
        ('ERA5 meteo\n+ local wind',    '1a6b8a'),
        ('Pre-process\nBC + forcing',   '2ba3a8'),
        ('DIMR run\nFM + SWAN',         '0d2137'),
        ('Post-process\nregrid + stats','2ba3a8'),
        ('Wetwise\nPortal bundle',      '3aaa6e'),
    ]
    bw, bh, by0 = 1.90, 1.10, 1.25
    for i, (lbl, col) in enumerate(steps):
        bx = 0.40 + i * 2.55
        add_rect(slide, bx, by0, bw, bh, rgb(col))
        add_text(slide, lbl, bx+0.08, by0+0.22, bw-0.16, 0.70,
                 size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i < len(steps)-1:
            add_text(slide, '→', bx+bw+0.08, by0+0.35, 0.50, 0.45,
                     size=22, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # Cadence box
    add_rect(slide, 0.35, 2.65, 7.5, 3.00, WHITE, line_color=MID, line_width=Pt(1.0))
    add_text(slide, 'Operational cadence', 0.55, 2.72, 7.1, 0.52,
             size=13, bold=True, color=MID)
    cadence = [
        'Daily 10-day forecast window (CMEMS anfc D+0→D+9)',
        'Wall time ~3.5 h / 9 days on simit-server\n'
        '  (8-core FM + 16-thread SWAN, Intel MPI)',
        'Output: map.nc + his.nc → regridded netCDF\n'
        '  + Plotly JSON bundles for portal',
        'Run chain: N−2 restart → 48 h run → drop 24 h\n'
        '  (validated approach, stable continuation)',
    ]
    bullet_box(slide, cadence, 0.55, 3.32, 7.0, 2.20,
               size=11, color=DARK, spacing_pt=14)

    # Data contract box
    add_rect(slide, 8.20, 2.65, 7.45, 3.00, WHITE, line_color=TEAL, line_width=Pt(1.0))
    add_text(slide, 'Data contract for portal', 8.40, 2.72, 7.0, 0.52,
             size=13, bold=True, color=TEAL)
    contract = [
        'manifest.json → run timestamp, variables,\n  valid date range',
        'Per-variable netCDF (50 m regular grid)\n  OR Plotly JSON for each tab',
        'Station time series: BocaNord, BocaSud,\n  AltaVilaEst, Marettimo',
        'Update trigger: webhook / S3 upload on\n  run completion',
    ]
    bullet_box(slide, contract, 8.40, 3.32, 7.0, 2.20,
               size=11, color=DARK, spacing_pt=14)

    # Milestones
    add_rect(slide, 0.35, 6.00, 15.30, 0.55, DARK)
    add_text(slide, 'MILESTONES', 0.50, 6.05, 5, 0.43,
             size=12, bold=True, color=WHITE)

    milestones = [
        ('Q3 2026', 'V2 Plotly bundle · variable roughness validation · seagrass ARL lagoon mask'),
        ('Q4 2026', 'Continuous 10-day operational forecast · portal live integration · v05 mesh'),
        ('2027',    'Paper 1 submission · D-Morph calibration with observed bathymetry · v05'),
    ]
    for i, (period, task) in enumerate(milestones):
        y = 6.65 + i * 0.72
        add_rect(slide, 0.35, y, 1.60, 0.55, MID if i<2 else TEAL)
        add_text(slide, period, 0.42, y+0.06, 1.46, 0.43,
                 size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, task, 2.10, y+0.06, 13.4, 0.43,
                 size=12, color=DARK)
    return slide


# ── Build and save ─────────────────────────────────────────────────────────────
def main():
    prs = new_prs()
    for fn in [slide_1_title, slide_2_model, slide_3_outputs,
               slide_4_panel, slide_5_demo, slide_6_roadmap]:
        fn(prs)
        print(f'  OK  {fn.__name__}')
    prs.save(OUT)
    print(f'\nSaved: {OUT}')


if __name__ == '__main__':
    main()
